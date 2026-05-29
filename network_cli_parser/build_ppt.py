"""Build PARSER_SOP.pptx — Network CLI Structured Parser SOP presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
C_DARK_BLUE   = RGBColor(0x1A, 0x35, 0x5E)   # slide title bg / header row
C_MID_BLUE    = RGBColor(0x1F, 0x5C, 0x99)   # accent / NEW badge bg
C_LIGHT_BLUE  = RGBColor(0xDB, 0xE9, 0xF7)   # table alt row
C_GREEN       = RGBColor(0x1A, 0x7A, 0x40)   # pass / positive
C_AMBER       = RGBColor(0xB3, 0x6B, 0x00)   # warn
C_RED         = RGBColor(0xB3, 0x1B, 0x1B)   # fail / negative
C_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
C_NEAR_WHITE  = RGBColor(0xF5, 0xF7, 0xFA)
C_CODE_BG     = RGBColor(0xEE, 0xF2, 0xF7)
C_GREY_TEXT   = RGBColor(0x44, 0x44, 0x44)
C_BODY_TEXT   = RGBColor(0x22, 0x22, 0x22)
C_PURPLE      = RGBColor(0x5B, 0x2D, 0x8E)   # branches colour

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_slide(prs, layout_idx=6):
    """Add a blank slide (layout 6 = blank)."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color and line_width:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    elif not line_color:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False, color=C_BODY_TEXT,
                align=PP_ALIGN.LEFT, font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_header_bar(slide, title, subtitle=None):
    """Navy header bar at top of slide."""
    bar = add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=C_DARK_BLUE)
    add_textbox(slide, Inches(0.35), Inches(0.1), Inches(12.5), Inches(0.65),
                title, font_size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.35), Inches(0.72), Inches(12.5), Inches(0.32),
                    subtitle, font_size=13, color=RGBColor(0xCC, 0xDD, 0xF0))
    return bar


def add_bullet_block(slide, x, y, w, h, items, font_size=14, color=C_BODY_TEXT,
                     bullet_char="▸ ", indent=0, line_spacing=1.15):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        prefix = "  " * indent + bullet_char
        run.text = prefix + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def add_code_block(slide, x, y, w, h, code_text, font_size=10.5):
    """Light-grey rounded code block."""
    bg = add_rect(slide, x, y, w, h, fill_color=C_CODE_BG,
                  line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.5))
    txBox = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.08),
                                     w - Inches(0.24), h - Inches(0.16))
    tf = txBox.text_frame
    tf.word_wrap = False
    first = True
    for line in code_text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Courier New"
        run.font.color.rgb = C_GREY_TEXT
    return bg, txBox


def add_table_slide(slide, x, y, w, col_widths, rows_data, header_color=C_DARK_BLUE,
                    alt_color=C_LIGHT_BLUE, font_size=12):
    """rows_data: list of lists (first row = header)."""
    rows = len(rows_data)
    cols = len(col_widths)
    row_h = Inches(0.38)
    h = row_h * rows
    tbl = slide.shapes.add_table(rows, cols, x, y, w, h).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = str(cell_text)
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.text = str(cell_text)
            run.font.size = Pt(font_size)
            run.font.name = "Calibri"
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
    return tbl


def add_badge(slide, x, y, label, bg=C_MID_BLUE, fg=C_WHITE, font_size=11):
    w = Inches(len(label) * 0.11 + 0.25)
    h = Inches(0.3)
    rect = add_rect(slide, x, y, w, h, fill_color=bg)
    add_textbox(slide, x + Inches(0.05), y + Inches(0.02), w - Inches(0.1), h,
                label, font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return w


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    slide = add_slide(prs)
    # Full-bleed background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK_BLUE)
    # Accent bar
    add_rect(slide, 0, Inches(4.9), SLIDE_W, Inches(0.08), fill_color=C_MID_BLUE)

    add_textbox(slide, Inches(0.7), Inches(1.6), Inches(11.5), Inches(1.2),
                "Network CLI Structured Parser", font_size=40, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.7), Inches(2.85), Inches(11.5), Inches(0.7),
                "Standard Operating Procedure", font_size=28, color=RGBColor(0xAA, 0xCC, 0xEE))
    add_textbox(slide, Inches(0.7), Inches(3.7), Inches(11.5), Inches(0.5),
                "Parsing · Health Checks · Compound Conditions · Reports · Collector",
                font_size=16, color=RGBColor(0x88, 0xAA, 0xCC))
    add_textbox(slide, Inches(0.7), Inches(5.5), Inches(11.5), Inches(0.4),
                "cisco_nxos  ·  cisco_ios  ·  TextFSM  ·  TTP  ·  YAML Checks",
                font_size=13, color=RGBColor(0x66, 0x88, 0xAA), italic=True)


def slide_toc(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Table of Contents")
    left_items = [
        "1. Overview & Architecture",
        "2. Running the Parser",
        "3. Directory Structure",
        "4. Status Codes",
        "5. Adding a New Command",
        "6. TextFSM Templates",
        "7. TTP Templates",
        "8. Pipe-Filtered Commands",
        "9. commands.yaml Reference",
        "10. Report Generator Overview",
    ]
    right_items = [
        "11. Health Check Path Syntax",
        "12. Conditions (14 total)",
        "13. match: all vs any",
        "14. Severity Levels",
        "15. ★ AND Conditions (NEW)",
        "16. ★ one_of / not_one_of (NEW)",
        "17. ★ Branches IF/ELIF/ELSE (NEW)",
        "18. ★ print Field (NEW)",
        "19. Per-Device Overrides & Baseline",
        "20. Delta · HTML Reports · Collector",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(5.8),
                     left_items, font_size=14, bullet_char="")
    add_bullet_block(slide, Inches(6.8), Inches(1.2), Inches(6.1), Inches(5.8),
                     right_items, font_size=14, bullet_char="")

    # Highlight NEW items
    add_rect(slide, Inches(6.65), Inches(4.2), Inches(6.4), Inches(1.7),
             fill_color=RGBColor(0xFF, 0xFA, 0xE5),
             line_color=RGBColor(0xFF, 0xCC, 0x00), line_width=Pt(1))
    add_textbox(slide, Inches(6.8), Inches(4.22), Inches(6.1), Inches(0.3),
                "★ = New features in this release", font_size=11,
                bold=True, color=C_AMBER)


def slide_overview(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Overview & Architecture",
                   "One JSON snapshot per device — every command, every output")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.35),
                "Parser Chain (per command):", font_size=14, bold=True, color=C_DARK_BLUE)

    chain_code = """\
commands.yaml lookup
    ├── raw_only       → skip parsing, preserve raw text
    ├── ntc            → NTC Templates (ntc-templates library)
    ├── custom         → Custom TextFSM template
    ├── ttp            → TTP template
    ├── hierarchical   → Python regex parser (multicast / VRF-nested)
    └── auto_discover  (unknown commands)
            ├── Step 1: NTC Templates (reconstructed command string)
            ├── Step 2: Convention TextFSM  {platform}_{cmd}.textfsm
            ├── Step 3: Convention TTP      {platform}_{cmd}.ttp
            └── no_template  (raw preserved, structured parsing skipped)"""

    add_code_block(slide, Inches(0.4), Inches(1.6), Inches(12.5), Inches(2.7),
                   chain_code, font_size=11)

    bullets = [
        "Every command found in the log appears in JSON output — even commands without a template",
        "Platforms supported: cisco_nxos, cisco_ios (extensible via commands.yaml)",
        "Auto-discovery requires no YAML edit — just drop a template file in the right folder",
        "JSON output: one file per device, containing metadata + commands dict",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(4.45), Inches(12.5), Inches(2.8),
                     bullets, font_size=13)


def slide_running(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Running the Parser  (main.py)")

    add_textbox(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.3),
                "Installation & Quick Start", font_size=14, bold=True, color=C_DARK_BLUE)

    code1 = """\
cd network_cli_parser
pip install -r requirements.txt

# Single .txt file
python main.py --input data/raw/N9K-CAMA-WAN-1_03-May-26.txt

# Entire directory (all .txt files)
python main.py --input data/raw/

# Custom output directory
python main.py --input data/raw/ --output /tmp/parsed/"""
    add_code_block(slide, Inches(0.4), Inches(1.6), Inches(7.5), Inches(2.45), code1)

    add_textbox(slide, Inches(8.3), Inches(1.25), Inches(4.8), Inches(0.3),
                "Output", font_size=14, bold=True, color=C_DARK_BLUE)
    bullets = [
        "One JSON file per input",
        "Written to data/json/ by default",
        "Filename: <hostname>_<date>.json",
        "Contains metadata + commands dict",
        "Raw output always preserved",
    ]
    add_bullet_block(slide, Inches(8.3), Inches(1.6), Inches(4.7), Inches(2.4),
                     bullets, font_size=13)

    add_textbox(slide, Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.3),
                "JSON Snapshot Structure", font_size=14, bold=True, color=C_DARK_BLUE)
    code2 = """\
{
  "metadata": { "hostname": "N9K-CAMA-WAN-1", "platform": "cisco_nxos",
                "collection_time": "2026-05-29T10:00:00" },
  "commands": {
    "show_version": { "status": "parsed", "raw": "...", "parsed": [{...}] },
    "show_ip_bgp_summary": { "status": "parsed", "raw": "...", "parsed": [{...}] }
  }
}"""
    add_code_block(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(2.5), code2)


def slide_status_codes(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Status Codes",
                   "Every command entry in the JSON carries one of these status values")

    rows = [
        ["Status", "Meaning", "Template Required?"],
        ["parsed", "Parser returned structured data rows", "Yes"],
        ["partial", "Parser ran but returned 0 rows (template matched nothing)", "Yes"],
        ["raw_only", "Explicitly registered as skip in commands.yaml", "No (intentional)"],
        ["failed", "Parser threw an exception", "Yes (broken template)"],
        ["no_template", "Not in registry; auto-discovery found nothing; raw preserved", "Build one"],
    ]
    col_widths = [Inches(1.8), Inches(7.6), Inches(2.8)]
    add_table_slide(slide, Inches(0.4), Inches(1.3), Inches(12.2), col_widths, rows)

    add_textbox(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.35),
                "How to handle no_template commands:", font_size=14, bold=True, color=C_DARK_BLUE)
    bullets = [
        "The raw output is always preserved — nothing is lost",
        "Write a TextFSM or TTP template, drop it in templates/ — re-run and it auto-discovers",
        "Or add an explicit entry to commands.yaml pointing at the new template",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(4.95), Inches(12.5), Inches(2.0),
                     bullets, font_size=13)


def slide_adding_commands(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Adding a New Command",
                   "Three options — pick the one that fits your workflow")

    # Option A
    add_badge(slide, Inches(0.4), Inches(1.2), "Option A — Explicit YAML registration", C_DARK_BLUE)
    code_a = """\
# commands.yaml
cisco_nxos:
  show ip pim neigh:
    parser: custom
    template: cisco_nxos_show_ip_pim_neigh"""
    add_code_block(slide, Inches(0.4), Inches(1.6), Inches(5.9), Inches(1.6), code_a, font_size=10.5)

    # Option B
    add_badge(slide, Inches(6.65), Inches(1.2), "Option B — Auto-discovery (no YAML edit)", C_MID_BLUE)
    code_b = """\
# Drop the file in the right folder:
templates/custom/routing/
  cisco_nxos_show_ip_pim_neigh.textfsm

# OR for TTP:
templates/ttp/routing/
  cisco_nxos_show_ip_pim_neigh.ttp

# Auto-discovered on next run — no YAML edit!"""
    add_code_block(slide, Inches(6.65), Inches(1.6), Inches(6.3), Inches(2.25), code_b, font_size=10.5)

    # Option C
    add_badge(slide, Inches(0.4), Inches(3.35), "Option C — Run first, build template later", C_GREEN)
    bullets_c = [
        "Run main.py with the new log — command appears as  status: no_template",
        "Raw output is preserved in JSON — no data is lost",
        "Build a template, drop it in templates/ — re-run and it auto-discovers",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(3.75), Inches(5.9), Inches(2.5),
                     bullets_c, font_size=12.5)

    # Naming convention
    add_textbox(slide, Inches(6.65), Inches(3.85), Inches(6.3), Inches(0.35),
                "Naming convention (auto-discovery):", font_size=13, bold=True, color=C_DARK_BLUE)
    code_naming = """\
Command: show ip pim neigh  (cisco_nxos)
Normalized key: show_ip_pim_neigh
TextFSM: cisco_nxos_show_ip_pim_neigh.textfsm
TTP:     cisco_nxos_show_ip_pim_neigh.ttp"""
    add_code_block(slide, Inches(6.65), Inches(4.25), Inches(6.3), Inches(1.65), code_naming, font_size=10.5)


def slide_textfsm(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "TextFSM Templates  (parser: custom)",
                   "State-machine parser — ideal for tabular show outputs")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Syntax", font_size=13, bold=True, color=C_DARK_BLUE)
    code_syntax = """\
Value FIELD_NAME (regex)
Value Filldown VRF (\\S+)   # Carries forward across rows

Start
  ^header line regex -> Continue
  ^${FIELD}\\s+${OTHER}  -> Record

EOF"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(2.1), code_syntax)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "Key directives", font_size=13, bold=True, color=C_DARK_BLUE)
    rows = [
        ["Directive", "Meaning"],
        ["Value X (re)", "Capture group, stored in field X"],
        ["Filldown X (re)", "Field carries forward until overwritten"],
        ["-> Record", "Emit current row to output"],
        ["-> Continue", "Re-process line in current state"],
        ["EOF", "Emit buffered row at end of input"],
    ]
    add_table_slide(slide, Inches(6.65), Inches(1.55), Inches(6.3),
                    [Inches(2.1), Inches(4.2)], rows, font_size=12)

    add_textbox(slide, Inches(0.4), Inches(3.75), Inches(12.5), Inches(0.3),
                "Example — Multi-VRF PIM neighbor table", font_size=13, bold=True, color=C_DARK_BLUE)
    code_ex = """\
Value Filldown VRF (\\S+)
Value NEIGHBOR (\\d+\\.\\d+\\.\\d+\\.\\d+)
Value UPTIME (\\S+)
Value STATE (\\S+)

Start
  ^PIM Neighbor.*VRF\\s+${VRF} -> Neighbors

Neighbors
  ^\\s*${NEIGHBOR}\\s+\\S+\\s+${UPTIME}\\s+${STATE} -> Record
  ^PIM Neighbor.*VRF\\s+${VRF} -> Continue.Record

EOF"""
    add_code_block(slide, Inches(0.4), Inches(4.1), Inches(12.5), Inches(3.0), code_ex, font_size=10.5)


def slide_ttp(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "TTP Templates  (parser: ttp)",
                   "Template Text Parser — ideal for hierarchical JSON output (VRF → neighbors)")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Flat neighbor table", font_size=13, bold=True, color=C_DARK_BLUE)
    code_flat = """\
<group name="neighbors">
{{ neighbor_id }}  {{ priority }}
  {{ state | re("\\\\S+/\\\\s*\\\\S*") }}
  {{ uptime }}  {{ address }}  {{ interface }}
</group>"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(1.85), code_flat)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "Hierarchical (VRF → neighbors)", font_size=13, bold=True, color=C_DARK_BLUE)
    code_hier = """\
<group name="vrfs">
BGP summary for VRF {{ vrf }}, af {{ afi | re(".+") }}
BGP router identifier {{ router_id }}

  <group name="neighbors">
  {{ neighbor | re("\\\\d+\\\\.\\\\d+\\\\.\\\\d+\\\\.\\\\d+") }}
    {{ remote_as }}  {{ updown }}  {{ state_pfx }}
  </group>

</group>"""
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(2.6), code_hier)

    add_textbox(slide, Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.3),
                "Result JSON structure (hierarchical)", font_size=13, bold=True, color=C_DARK_BLUE)
    code_json = """\
{ "vrfs": [{ "vrf": "default", "afi": "IPv4 Unicast", "router_id": "10.0.0.1",
             "neighbors": [{ "neighbor": "10.0.0.2", "remote_as": "65002",
                             "updown": "2w3d", "state_pfx": "100" }] }] }"""
    add_code_block(slide, Inches(0.4), Inches(3.85), Inches(12.5), Inches(1.1), code_json, font_size=10.5)

    bullets = [
        "Use re() when field contains spaces  (e.g. FULL/ - would break default whitespace split)",
        "Use an IP-anchored regex on the first field to avoid accidentally matching header lines",
        "Naming: {platform}_{normalized_cmd}.ttp  →  cisco_nxos_show_ip_bgp_summary_vrf_all.ttp",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(5.1), Inches(12.5), Inches(2.0),
                     bullets, font_size=13)


def slide_pipe_filter(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Pipe-Filtered Commands  ( | include / | grep … )",
                   "Piped output has a different structure — it gets its own key and template")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.3),
                "Normalization behaviour", font_size=13, bold=True, color=C_DARK_BLUE)
    rows = [
        ["Raw command", "Normalized key"],
        ["show ip bgp summary", "show_ip_bgp_summary"],
        ["show ip bgp summary | include 65001", "show_ip_bgp_summary_include"],
        ["show ip bgp summary | exclude Connected", "show_ip_bgp_summary_exclude"],
        ["show ip route | begin 10.0.0", "show_ip_route_begin"],
        ["show ip route | grep 10.0.0", "show_ip_route_grep"],
        ["show ip route | count", "show_ip_route_count"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.55), Inches(12.2),
                    [Inches(6.8), Inches(5.4)], rows, font_size=12)

    add_textbox(slide, Inches(0.4), Inches(4.0), Inches(6.0), Inches(0.3),
                "Template file naming (auto-discovered):", font_size=13, bold=True, color=C_DARK_BLUE)
    code = """\
# TextFSM (drop in templates/custom/):
cisco_nxos_show_ip_bgp_summary_include.textfsm

# TTP (drop in templates/ttp/):
cisco_nxos_show_ip_bgp_summary_include.ttp

# No YAML edit needed — auto-discovered"""
    add_code_block(slide, Inches(0.4), Inches(4.35), Inches(6.0), Inches(2.4), code, font_size=10.5)

    add_textbox(slide, Inches(6.65), Inches(4.0), Inches(6.3), Inches(0.3),
                "Explicit YAML registration:", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
cisco_nxos:
  show ip bgp summary | include:
    parser: custom
    template: cisco_nxos_show_ip_bgp_summary_include

# Note: use the literal '| include' as the key
# The mapper normalizes it at load time"""
    add_code_block(slide, Inches(6.65), Inches(4.35), Inches(6.3), Inches(2.4), code2, font_size=10.5)


def slide_report_overview(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Report Generator  (report.py)",
                   "Operates on JSON snapshots; supports JSON + self-contained HTML output")

    rows = [
        ["Subcommand", "Purpose", "Output"],
        ["delta", "Field diff between two snapshots", "JSON / HTML"],
        ["delta-all", "Per-device delta across two directories", "HTML + index.html"],
        ["health", "YAML check file vs single snapshot", "JSON / HTML (exit 1 on critical fail)"],
        ["health-all", "Health checks across all snapshots in a dir", "HTML + index.html"],
        ["baseline", "Auto-generate starter check YAML from current values", "YAML file"],
        ["collect", "Collect via SSH (online) or process .txt dumps (offline)", "JSON snapshots"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.3), Inches(12.5),
                    [Inches(1.8), Inches(6.8), Inches(3.6)], rows, font_size=12)

    add_textbox(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.3),
                "Common invocations", font_size=13, bold=True, color=C_DARK_BLUE)
    code = """\
python report.py health --snapshot data/json/N9K.json --checks checks/base.yaml
python report.py health --snapshot data/json/N9K.json --checks checks/base.yaml --output health.html
python report.py health-all --dir data/json/ --default-checks checks/base.yaml --output-dir reports/
python report.py delta --before before/N9K.json --after after/N9K.json --output delta.html
python report.py delta-all --before-dir snapshots/03-May/ --after-dir snapshots/04-May/ --output-dir reports/
python report.py baseline --snapshot data/json/N9K.json --output checks/baseline.yaml
python report.py collect --from-dir data/raw/ --output-dir data/json/"""
    add_code_block(slide, Inches(0.4), Inches(4.9), Inches(12.5), Inches(2.4), code, font_size=10)


def slide_path_syntax(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Health Check — Path Syntax",
                   "Paths navigate the parsed JSON structure of a command")

    rows = [
        ["Token", "Meaning"],
        ["[0]", "Index into a list  →  first element"],
        ["[*]", "Expand ALL items of a list OR all values of a dict"],
        [".field", "Access a dict key  (dot is optional at start of path)"],
        ["field", "Dict key at the start of a path (no leading dot needed)"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.25), Inches(12.2),
                    [Inches(1.8), Inches(10.4)], rows)

    add_textbox(slide, Inches(0.4), Inches(3.15), Inches(12.5), Inches(0.3),
                "Examples", font_size=13, bold=True, color=C_DARK_BLUE)
    rows2 = [
        ["Path", "What it accesses"],
        ["[0].os", "First element of a list, os field"],
        ["[*].status", "status field from every row in a list"],
        ["vrfs[*].summary.total_routes", "total_routes from summary of every VRF in a dict"],
        ["neighbors[*].state", "state from every neighbor in a list"],
        ["goid[*].peer[*].ike_status", "ike_status from every peer inside every GOID"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(3.5), Inches(12.2),
                    [Inches(4.2), Inches(8.0)], rows2)

    add_textbox(slide, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.5),
                "When [*] expands multiple values → condition is checked against each one  "
                "(all must pass by default; use  match: any  to require at least one)",
                font_size=12, color=C_GREY_TEXT, italic=True)


def slide_conditions(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Conditions — 14 Built-in Operators",
                   "Use the condition field in a check definition")

    rows = [
        ["Condition", "Passes when (actual vs expected)"],
        ["eq", "actual == expected"],
        ["ne", "actual != expected"],
        ["gt / gte", "actual >  /  >=  expected (numeric)"],
        ["lt / lte", "actual <  /  <=  expected (numeric)"],
        ["contains", "str(actual) contains expected substring"],
        ["not_contains", "str(actual) does NOT contain expected substring"],
        ["matches", "str(actual) matches expected regex  (re.search)"],
        ["one_of  ★", "actual is in the expected list  →  value: [a, b, c]"],
        ["not_one_of  ★", "actual is NOT in the expected list"],
        ["duration_gt / duration_gte", "parsed duration  >  /  >=  expected duration"],
        ["duration_lt / duration_lte", "parsed duration  <  /  <=  expected duration"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.25), Inches(12.5),
                    [Inches(3.0), Inches(9.5)], rows, font_size=12)

    add_textbox(slide, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.5),
                "★ one_of / not_one_of are new — see slide 16 for full examples  |  "
                "Duration conditions parse strings like '5w2d', '00:03:42', 'never'",
                font_size=12, color=C_GREY_TEXT, italic=True)


def slide_duration(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Duration Conditions",
                   "Compare uptime / timer strings without manual string parsing")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Supported input formats", font_size=13, bold=True, color=C_DARK_BLUE)
    rows = [
        ["Input", "Interpreted as"],
        ["5w2d", "5 weeks + 2 days"],
        ["2d03h", "2 days + 3 hours"],
        ["1y3w2d4h5m6s", "full decomposition"],
        ["00:03:42", "HH:MM:SS"],
        ["15:30", "MM:SS"],
        ["5week2day", "spelled-out units"],
        ["never / n/a / -", "0 seconds"],
        ["42", "42 seconds (bare integer)"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.55), Inches(6.0),
                    [Inches(2.6), Inches(3.4)], rows, font_size=12)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "Example checks", font_size=13, bold=True, color=C_DARK_BLUE)
    code = """\
# BGP session up at least 2 days
- name: "BGP session stable"
  command: show_ip_bgp_summary_vrf_all
  path: "vrfs[*].neighbors[*].updown"
  condition: duration_gte
  value: "2d"

# OSPF neighbor up more than 1 week
- name: "OSPF neighbor stable"
  command: show_ip_ospf_neighbors
  path: "neighbors[*].uptime"
  condition: duration_gt
  value: "1w"

# PIM neighbor recently reset (< 1h)
- name: "PIM neighbor fresh"
  command: show_ip_pim_neigh
  path: "[*].UPTIME"
  condition: duration_lt
  value: "1h" """
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(5.0), code)

    bullets = [
        "Both  actual  (from snapshot) and  value  (from YAML) go through the same parser",
        "Use any supported format in either field — they don't need to match",
        "0 is NOT > 0 — 'never' fails  duration_gt: '0s'",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(5.35), Inches(6.0), Inches(2.0),
                     bullets, font_size=12)


def slide_match_severity(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "match  and  severity  Fields",
                   "Control how multi-value expansion is evaluated + how failures affect CI")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "match: all / any", font_size=14, bold=True, color=C_DARK_BLUE)
    rows_match = [
        ["match value", "Default?", "Passes when"],
        ["all", "Yes", "Every expanded value satisfies the condition"],
        ["any", "No", "At least one expanded value satisfies the condition"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.55), Inches(6.0),
                    [Inches(1.5), Inches(1.3), Inches(3.2)], rows_match, font_size=12)

    code_match = """\
# At least one IKE peer is Established
- name: "IKE peer alive"
  command: show_crypto_gkm_ks_coop_detail
  path: "goid[*].peer[*].ike_status"
  condition: matches
  value: "Established"
  match: any   # pass if ANY peer is Established"""
    add_code_block(slide, Inches(0.4), Inches(2.9), Inches(6.0), Inches(2.6), code_match, font_size=10.5)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "severity: critical / warn / info", font_size=14, bold=True, color=C_DARK_BLUE)
    rows_sev = [
        ["severity", "Default?", "Exit code on failure"],
        ["critical", "Yes", "exit 1  — blocks CI"],
        ["warn", "No", "exit 0  — shown in report only"],
        ["info", "No", "exit 0  — informational"],
    ]
    add_table_slide(slide, Inches(6.65), Inches(1.55), Inches(6.3),
                    [Inches(1.5), Inches(1.3), Inches(3.5)], rows_sev, font_size=12)

    code_sev = """\
# MTU mismatch is a warning, not a blocker
- name: "MTU is jumbo"
  command: show_interfaces
  path: "[*].mtu"
  condition: eq
  value: 9216
  severity: warn

# Baseline checks default to warn so they
# never block CI until explicitly promoted"""
    add_code_block(slide, Inches(6.65), Inches(2.9), Inches(6.3), Inches(2.6), code_sev, font_size=10.5)

    add_textbox(slide, Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.3),
                "Both match and severity are optional — omitting them uses the defaults (all / critical)",
                font_size=12, color=C_GREY_TEXT, italic=True)


def slide_and_conditions(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "★  AND Conditions  (New)  — conditions list",
                   "All conditions in the list must pass for each resolved value")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.3),
                "Use case: range checks — value must be > min AND < max",
                font_size=14, bold=True, color=C_DARK_BLUE)

    code = """\
# BGP prefix count must be between 1 and 500,000 (both bounds must hold)
- name: "BGP prefix count in valid range"
  command: show_ip_bgp_summary
  path: "neighbors[*].prefixes_received"
  conditions:                        # ← list replaces single condition/value
    - condition: gte
      value: 1
    - condition: lte
      value: 500000
  severity: critical

# Works with match: any too
- name: "At least one interface has valid MTU"
  command: show_interfaces
  path: "[*].mtu"
  conditions:
    - condition: gte
      value: 1500
    - condition: lte
      value: 9216
  match: any"""
    add_code_block(slide, Inches(0.4), Inches(1.6), Inches(7.5), Inches(4.6), code)

    add_textbox(slide, Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.3),
                "How it works", font_size=13, bold=True, color=C_DARK_BLUE)
    bullets = [
        "Each item in conditions has its own condition + value",
        "An item passes only if ALL conditions pass (AND logic)",
        "On failure, every failing sub-condition is listed separately",
        "Compatible with match: all (default) and match: any",
        "Cannot combine conditions with branches in the same check",
    ]
    add_bullet_block(slide, Inches(8.2), Inches(1.6), Inches(4.8), Inches(3.0),
                     bullets, font_size=12.5)

    add_textbox(slide, Inches(8.2), Inches(4.7), Inches(4.8), Inches(0.3),
                "Failure output (terminal)", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
✗ BGP prefix count   [fail]
  path:    neighbors[0].prefixes_received
  actual:  600000
  reason:  600000 not <= 500000"""
    add_code_block(slide, Inches(8.2), Inches(5.05), Inches(4.8), Inches(1.65), code2, font_size=10)


def slide_one_of(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "★  one_of  and  not_one_of  (New)",
                   "Value must be in (or not in) a list of allowed values")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "one_of — value must be in list", font_size=13, bold=True, color=C_DARK_BLUE)
    code1 = """\
# Interface protocol must be 'up' or 'connected'
- name: "Interface protocol valid"
  command: show_interfaces
  path: "[*].line_protocol"
  condition: one_of
  value: ["up", "connected"]

# BGP state must be Established or Active
- name: "BGP peer state acceptable"
  command: show_ip_bgp_summary
  path: "neighbors[*].state_pfx"
  condition: one_of
  value: ["Established", "Active", "Idle"]
  severity: warn"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(3.5), code1)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "not_one_of — value must NOT be in list", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
# No interface description should be DECOM
- name: "No decom interface"
  command: show_interfaces
  path: "[*].description"
  condition: not_one_of
  value: ["DECOM", "DECOMMISSIONED", "DELETE"]

# Spanning tree state — block bad states
- name: "No STP blocking"
  command: show_spanning_tree
  path: "[*].port_state"
  condition: not_one_of
  value: ["BLK", "discarding"]
  severity: warn"""
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(3.5), code2)

    add_textbox(slide, Inches(0.4), Inches(5.15), Inches(12.5), Inches(0.3),
                "Key points", font_size=13, bold=True, color=C_DARK_BLUE)
    bullets = [
        "value must be a YAML list  [item1, item2, ...]  — a plain string value causes an error",
        "Works with match: all (default) or match: any — e.g. require at least one interface to be up",
        "Combines with severity: warn/info to make non-blocking checks",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(5.5), Inches(12.5), Inches(1.6),
                     bullets, font_size=13)


def slide_branches(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "★  Conditional Branches  IF / ELIF / ELSE  (New)",
                   "Apply different checks depending on a field value in each row")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(7.5), Inches(0.3),
                "Syntax", font_size=13, bold=True, color=C_DARK_BLUE)
    code = """\
- name: "MTU matches interface type"
  command: show_interfaces
  path: "[*]"               # must resolve to dict rows
  branches:
    - when:                 # IF
        field: type
        condition: eq
        value: loopback
      then:
        field: mtu
        condition: eq
        value: 65535
    - when:                 # ELIF
        field: type
        condition: contains
        value: Ethernet
      then:
        field: mtu
        condition: eq
        value: 9216
    - default:              # ELSE
        field: mtu
        condition: gte
        value: 1500"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(7.5), Inches(5.55), code)

    add_textbox(slide, Inches(8.2), Inches(1.2), Inches(4.8), Inches(0.3),
                "How it works", font_size=13, bold=True, color=C_DARK_BLUE)
    bullets = [
        "path must resolve to dict rows",
        "when: field/condition/value is the IF guard",
        "Branches are evaluated in order — first match wins",
        "default fires if no when matched (ELSE)",
        "No match + no default → row passes vacuously",
        "Non-dict row → fail with type error",
        "Compatible with match: all / any",
        "Compatible with print field",
    ]
    add_bullet_block(slide, Inches(8.2), Inches(1.55), Inches(4.8), Inches(3.5),
                     bullets, font_size=12.5)

    add_textbox(slide, Inches(8.2), Inches(5.05), Inches(4.8), Inches(0.3),
                "then / default structure", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
then:            # or default:
  field: mtu     # field to evaluate in the row
  condition: eq  # any standard condition
  value: 9216    # expected value"""
    add_code_block(slide, Inches(8.2), Inches(5.4), Inches(4.8), Inches(1.65), code2, font_size=10.5)


def slide_print_field(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "★  print Field  (New)  — Surface Values in Output",
                   "Display resolved values on every check — pass or fail")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Two forms", font_size=13, bold=True, color=C_DARK_BLUE)

    code1 = """\
# print: true  — raw path → value
- name: "Current OS version"
  command: show_version
  path: "[0].version"
  condition: matches
  value: "17\\\\..*"
  print: true
# Terminal output:
#  ✓ Current OS version   [pass]
#    [0].version -> '17.6.1a'

# print: "template"  — compose a sentence
- name: "BGP prefix count"
  command: show_ip_bgp_summary
  path: "neighbors[*].prefixes_received"
  conditions:
    - {condition: gte, value: 1}
    - {condition: lte, value: 500000}
  print: "Peer {path} has {value} prefixes"
# Terminal output per peer:
#  Peer neighbors[0].prefixes_received has 12453 prefixes"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.5), Inches(5.2), code1)

    add_textbox(slide, Inches(7.2), Inches(1.2), Inches(5.8), Inches(0.3),
                "Template placeholders", font_size=13, bold=True, color=C_DARK_BLUE)
    rows = [
        ["Placeholder", "Replaced with"],
        ["{value}", "The resolved actual value (str)"],
        ["{path}", "The resolved path string  e.g. neighbors[0].state"],
    ]
    add_table_slide(slide, Inches(7.2), Inches(1.55), Inches(5.8),
                    [Inches(2.0), Inches(3.8)], rows, font_size=12)

    add_textbox(slide, Inches(7.2), Inches(2.65), Inches(5.8), Inches(0.3),
                "Works with all check types", font_size=13, bold=True, color=C_DARK_BLUE)
    bullets2 = [
        "Simple checks  (condition + value)",
        "AND checks  (conditions list)",
        "Branches  (when/then/default)",
        "Shown on BOTH pass and fail",
        "One printed line per resolved item",
        "Also rendered in HTML report",
    ]
    add_bullet_block(slide, Inches(7.2), Inches(3.0), Inches(5.8), Inches(2.5),
                     bullets2, font_size=12.5)

    code_branches = """\
# print with branches
- name: "MTU by type"
  command: show_interfaces
  path: "[*]"
  print: "Interface {path} MTU is {value}"
  branches:
    - when: {field: type, condition: eq, value: loopback}
      then: {field: mtu, condition: eq, value: 65535}
    - default: {field: mtu, condition: gte, value: 1500}"""
    add_code_block(slide, Inches(7.2), Inches(5.35), Inches(5.8), Inches(1.8), code_branches, font_size=10)


def slide_compound_full_example(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Compound Conditions — Full Example",
                   "Combining conditions, branches, print, match, and severity in one check file")

    code = """\
checks:
  # AND conditions — range check
  - name: "BGP prefix count in range"
    command: show_ip_bgp_summary
    path: "neighbors[*].prefixes_received"
    conditions:
      - condition: gte
        value: 1
      - condition: lte
        value: 500000
    print: "Peer {path} has {value} prefixes"
    severity: critical

  # one_of — valid protocol states
  - name: "Interface protocol valid"
    command: show_interfaces
    path: "[*].line_protocol"
    condition: one_of
    value: ["up", "connected"]
    match: all
    severity: critical

  # not_one_of — no decom interfaces
  - name: "No DECOM interface"
    command: show_interfaces
    path: "[*].description"
    condition: not_one_of
    value: ["DECOM", "DECOMMISSIONED"]
    severity: warn

  # Branches IF/ELIF/ELSE — MTU by interface type
  - name: "MTU matches interface type"
    command: show_interfaces
    path: "[*]"
    print: "Interface {path} MTU={value}"
    branches:
      - when: {field: type, condition: eq, value: loopback}
        then: {field: mtu, condition: eq, value: 65535}
      - when: {field: type, condition: contains, value: Ethernet}
        then: {field: mtu, condition: eq, value: 9216}
      - default: {field: mtu, condition: gte, value: 1500}"""
    add_code_block(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(6.0), code, font_size=10)


def slide_per_device(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Per-Device Check Overrides  &  Baseline Generation",
                   "Two-tier check system: default checks + per-device overrides")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Per-device override file", font_size=13, bold=True, color=C_DARK_BLUE)
    code1 = """\
# checks/devices/N9K-CAMA-WAN-1.yaml
checks:
  # OVERRIDE: replaces default check (same name)
  - name: "VPC peer link status is up"
    command: show_vpc_brief
    path: "[*].status"
    condition: not_contains
    value: "peer-link down"  # more permissive

  # ADDITION: only for this device
  - name: "Hostname matches expected"
    command: show_version
    path: "[0].hostname"
    condition: eq
    value: "N9K-CAMA-WAN-1"

# Merge rule: device check WINS on name clash
# All other default checks remain unchanged"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(4.15), code1)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "CLI usage", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
# Single device with override
python report.py health \\
  --snapshot      data/json/N9K.json \\
  --checks        checks/base.yaml \\
  --device-checks checks/devices/N9K.yaml

# All devices — auto-loads {hostname}.yaml
python report.py health-all \\
  --dir               data/json/ \\
  --default-checks    checks/base.yaml \\
  --device-checks-dir checks/devices/ \\
  --output-dir        reports/health/"""
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(2.7), code2)

    add_textbox(slide, Inches(6.65), Inches(4.35), Inches(6.3), Inches(0.3),
                "Auto-baseline generation", font_size=13, bold=True, color=C_DARK_BLUE)
    code3 = """\
python report.py baseline \\
  --snapshot data/json/N9K.json \\
  --output   checks/baseline_N9K.yaml

# Generates eq checks for every scalar field
# Default severity: warn (won't block CI)
# Review and promote critical ones manually"""
    add_code_block(slide, Inches(6.65), Inches(4.7), Inches(6.3), Inches(2.1), code3)

    bullets = [
        "File naming: checks/devices/{hostname}.yaml — loaded automatically in health-all",
        "Devices with no matching file use only the default checks",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(5.8), Inches(6.0), Inches(1.3),
                     bullets, font_size=12)


def slide_delta(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Delta Report — Field-Level Diff Between Snapshots",
                   "Natural-key row matching — position changes are not reported as diffs")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Natural key priority", font_size=13, bold=True, color=C_DARK_BLUE)
    rows = [
        ["Natural key tried (in order)", "Example"],
        ["INTERFACE / interface / port", "Interface tables"],
        ["NEIGHBOR / neighbor / neighbor_id", "BGP / OSPF neighbor tables"],
        ["NETWORK / PREFIX / network", "Route tables"],
        ["VLAN / vlan_id", "VLAN tables"],
        ["Index (fallback)", "Any other list"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.55), Inches(6.0),
                    [Inches(3.2), Inches(2.8)], rows, font_size=12)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "JSON output structure", font_size=13, bold=True, color=C_DARK_BLUE)
    code = """\
{
  "metadata": {
    "before": {"hostname": "...", "collection_time": "..."},
    "after":  {"hostname": "...", "collection_time": "..."}
  },
  "summary": {
    "commands_added":     ["cmd_a"],
    "commands_removed":   [],
    "commands_changed":   ["show_vpc_brief"],
    "commands_unchanged": ["show_version", "..."]
  },
  "changes": {
    "show_vpc_brief": {
      "diffs": [
        {
          "path":   "parsed[port=Po10].status",
          "before": "up",
          "after":  "down"
        }
      ]
    }
  }
}"""
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(4.4), code)

    bullets = [
        "Row moving position in a list is NOT reported — only genuine field value changes are",
        "Commands present in only one snapshot appear as added / removed",
        "delta-all matches devices across directories by metadata.hostname",
    ]
    add_bullet_block(slide, Inches(0.4), Inches(5.0), Inches(6.0), Inches(2.1),
                     bullets, font_size=12.5)


def slide_html_reports(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "HTML Reports — Self-Contained, No Internet Required",
                   "All CSS and JS embedded; single .html file per device + index.html for multi-device")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(0.3),
                "Health report layout", font_size=13, bold=True, color=C_DARK_BLUE)
    rows_h = [
        ["Section", "Description"],
        ["Summary cards", "Total / Passed / Failed / Errors at a glance"],
        ["Check results", "Color-coded card per check; failures show path, actual, reason"],
        ["Severity badges", "WARN / INFO badges on non-critical failures"],
        ["Printed values", "Rendered per-item printed list (pass and fail)"],
        ["Raw outputs", "Every command's raw CLI text; Expand All / Collapse All"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(1.55), Inches(5.9),
                    [Inches(2.1), Inches(3.8)], rows_h, font_size=11)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "Delta report layout", font_size=13, bold=True, color=C_DARK_BLUE)
    rows_d = [
        ["Section", "Description"],
        ["Metadata", "Before → After hostname and timestamp side by side"],
        ["Summary cards", "Added / Removed / Changed / Unchanged"],
        ["Added / removed", "Command names as pills"],
        ["Changed commands", "Diff table: path / before (red) / after (green)"],
        ["Raw outputs", "Changed: before + after side by side"],
    ]
    add_table_slide(slide, Inches(6.65), Inches(1.55), Inches(6.3),
                    [Inches(2.1), Inches(4.2)], rows_d, font_size=11)

    add_textbox(slide, Inches(0.4), Inches(4.55), Inches(12.5), Inches(0.3),
                "Index pages (health-all and delta-all)", font_size=13, bold=True, color=C_DARK_BLUE)
    rows_idx = [
        ["Index page", "Columns", "Row color"],
        ["health-all index", "Hostname / Timestamp / Total / Passed / Failed / Errors / Report",
         "Muted if all-pass; red if any failure"],
        ["delta-all index", "Hostname / Status / Added / Removed / Changed / Unchanged / Report",
         "Muted if clean; faded if unmatched"],
    ]
    add_table_slide(slide, Inches(0.4), Inches(4.9), Inches(12.5),
                    [Inches(2.2), Inches(7.3), Inches(3.0)], rows_idx, font_size=11)

    add_textbox(slide, Inches(0.4), Inches(6.65), Inches(12.5), Inches(0.45),
                "Devices appearing in only one directory (delta-all) are listed as UNMATCHED with no report link.",
                font_size=12, color=C_GREY_TEXT, italic=True)


def slide_collector(prs):
    slide = add_slide(prs)
    add_header_bar(slide, "Snapshot Collector  (report.py collect)",
                   "Full pipeline: from live device to JSON snapshot — no main.py call needed")

    add_textbox(slide, Inches(0.4), Inches(1.2), Inches(6.0), Inches(0.3),
                "Offline mode (no SSH needed)", font_size=13, bold=True, color=C_DARK_BLUE)
    code1 = """\
# Process existing .txt CLI dumps
python report.py collect \\
  --from-dir   data/raw/ \\
  --output-dir data/json/

# Same parsing pipeline as main.py"""
    add_code_block(slide, Inches(0.4), Inches(1.55), Inches(6.0), Inches(1.9), code1)

    add_textbox(slide, Inches(0.4), Inches(3.55), Inches(6.0), Inches(0.3),
                "SSH mode (requires  pip install netmiko)", font_size=13, bold=True, color=C_DARK_BLUE)
    code2 = """\
python report.py collect \\
  --devices    devices.yaml \\
  --raw-dir    data/raw/      # optional .txt dumps
  --output-dir data/json/ \\
  --password   <pw>"""
    add_code_block(slide, Inches(0.4), Inches(3.9), Inches(6.0), Inches(1.85), code2)

    add_textbox(slide, Inches(6.65), Inches(1.2), Inches(6.3), Inches(0.3),
                "devices.yaml format", font_size=13, bold=True, color=C_DARK_BLUE)
    code3 = """\
defaults:
  username: admin
  password: ""      # leave blank, use --password
  timeout: 30

devices:
  - hostname: N9K-CAMA-WAN-1
    host: 192.168.1.1
    platform: cisco_nxos

  - hostname: IOS-RTR-1
    host: 10.0.0.1
    platform: cisco_ios
    username: ops        # override defaults

  - hostname: N9K-CAMA-WAN-2
    host: 192.168.1.2
    platform: cisco_nxos
    commands:            # explicit list (overrides registry)
      - show version
      - show ip bgp summary vrf all"""
    add_code_block(slide, Inches(6.65), Inches(1.55), Inches(6.3), Inches(3.7), code3)

    add_textbox(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.3),
                "End-to-end workflow:", font_size=13, bold=True, color=C_DARK_BLUE)
    code4 = """\
python report.py collect --devices devices.yaml --output-dir data/json/
python report.py delta-all --before-dir data/json-prev/ --after-dir data/json/ --output-dir reports/delta/
python report.py health-all --dir data/json/ --default-checks checks/base.yaml --output-dir reports/health/"""
    add_code_block(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(1.05), code4, font_size=10.5)


def slide_closing(prs):
    slide = add_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=C_DARK_BLUE)
    add_rect(slide, 0, Inches(3.5), SLIDE_W, Inches(0.08), fill_color=C_MID_BLUE)

    add_textbox(slide, Inches(1.0), Inches(1.2), Inches(11.0), Inches(1.0),
                "Network CLI Structured Parser", font_size=34, bold=True, color=C_WHITE,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.0), Inches(2.3), Inches(11.0), Inches(0.7),
                "Quick Reference Card", font_size=22, color=RGBColor(0xAA, 0xCC, 0xEE),
                align=PP_ALIGN.CENTER)

    summary = [
        "main.py  →  Parse .txt dumps to JSON  (data/json/)",
        "report.py collect  →  SSH or offline batch collection",
        "report.py health  →  YAML checks against snapshot (exit 1 on critical fail)",
        "report.py delta   →  Field-level diff between two snapshots",
        "report.py baseline →  Auto-generate starter check YAML",
        "Compound checks: conditions (AND) · one_of / not_one_of · branches (IF/ELIF/ELSE)",
        "print field: surface resolved values in terminal + HTML on pass and fail",
        "match: any — at least one value must pass;  severity: warn/info — no CI block",
    ]
    add_bullet_block(slide, Inches(1.0), Inches(3.75), Inches(11.0), Inches(3.3),
                     summary, font_size=14, color=RGBColor(0xCC, 0xDD, 0xF0))

    add_textbox(slide, Inches(1.0), Inches(7.1), Inches(11.0), Inches(0.3),
                "Full reference: PARSER_SOP.md  |  Starter checks: checks/example_health_checks.yaml",
                font_size=12, color=RGBColor(0x66, 0x88, 0xAA), align=PP_ALIGN.CENTER, italic=True)


# ── Main ──────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_toc(prs)
    slide_overview(prs)
    slide_running(prs)
    slide_status_codes(prs)
    slide_adding_commands(prs)
    slide_textfsm(prs)
    slide_ttp(prs)
    slide_pipe_filter(prs)
    slide_report_overview(prs)
    slide_path_syntax(prs)
    slide_conditions(prs)
    slide_duration(prs)
    slide_match_severity(prs)
    slide_and_conditions(prs)
    slide_one_of(prs)
    slide_branches(prs)
    slide_print_field(prs)
    slide_compound_full_example(prs)
    slide_per_device(prs)
    slide_delta(prs)
    slide_html_reports(prs)
    slide_collector(prs)
    slide_closing(prs)

    out = "PARSER_SOP.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
